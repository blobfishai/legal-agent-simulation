#!/usr/bin/env python3
"""Generate Harbor-format tasks (github.com/harbor-framework/harbor) from the
canonical world document.

Each of the world's tasks becomes one Harbor task directory:

    dist/harbor/tasks/task_XXX/
      instruction.md                the task prompt + interaction contract
      task.toml                     schema 1.4 config + provenance metadata
      environment/Dockerfile        agent container (pinned python + `tool` CLI)
      environment/tool              firm-systems CLI (JSON-RPC over the shim)
      environment/docker-compose.yaml  adds the shared `world` service
      tests/test.sh                 POST /verify -> /logs/verifier/reward.json
      solution/solve.sh             token-gated POST /solve (oracle reference walk)

The world itself (runtime + world doc + product contracts + shim) is ONE shared
Docker image (dist/harbor/world-image/, built with --build-image); per-task
compose files select the task via the TASK_ID env var. The agent container
never contains world.json, so verifier code and reference walks are not
readable by the agent.

Usage:
  python3 harbor/generate.py [--world world/blobfish/world-v21.json]
                             [--contracts mcp/v5/contracts]
                             [--out dist/harbor] [--tasks task_003,task_010]
                             [--build-image] [--image-tag legal-agent-sim-world:v21]
"""
from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import secrets
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path, PurePosixPath

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)

DISCLAIMER = ("Simulation only - every matter, client, document, attorney, "
              "and figure is synthetic test data.")
WORLD_RUNTIME_FILES = (
    "server.py",
    "oracle.py",
    "v2runtime.py",
    "v3dialects.py",
    "evidence.py",
    "paging.py",
    "wire_errors.py",
    "product_workflows.py",
    "query_dsl.py",
    "v21_verifier_runtime.py",
)
WORLD_IMAGE_TEMPLATE_FILES = ("shim.py", "start.sh", "Dockerfile")
SAFE_COMPONENT = re.compile(r"^[A-Za-z0-9][A-Za-z0-9._-]{0,127}$")
_VALIDATED_SOURCE_TREES: set[Path] = set()


def safe_component(value: object, label: str) -> str:
    if not isinstance(value, str) or not SAFE_COMPONENT.fullmatch(value):
        raise RuntimeError(f"{label} must be one safe filesystem component")
    return value


def resolve_output_root(value: str | os.PathLike[str]) -> Path:
    """Confine generated/replaced trees to a non-symlinked child of dist/."""
    allowed = (Path(ROOT) / "dist").resolve()
    lexical = Path(os.path.abspath(os.fspath(value)))
    if lexical == allowed or allowed not in lexical.parents:
        raise RuntimeError(f"--out must be a child of {allowed}")
    current = allowed
    for part in lexical.relative_to(allowed).parts:
        current /= part
        if current.is_symlink():
            raise RuntimeError(f"--out may not contain a symlink component: {current}")
    resolved = lexical.resolve()
    if resolved == allowed or allowed not in resolved.parents:
        raise RuntimeError(f"--out resolves outside {allowed}: {resolved}")
    return resolved


def reset_generated_directory(path: Path, label: str) -> None:
    if path.is_symlink():
        raise RuntimeError(f"refusing to replace symlinked {label}: {path}")
    if path.exists():
        if not path.is_dir():
            raise RuntimeError(f"refusing to replace non-directory {label}: {path}")
        shutil.rmtree(path)
    path.mkdir(parents=True)


def validate_project_source_tree(path: Path, label: str) -> Path:
    """Reject links, special files, and world-directed reads outside the repo."""
    repository = Path(ROOT).resolve()
    lexical = Path(os.path.abspath(path))
    if lexical == repository or repository not in lexical.parents:
        raise RuntimeError(f"{label} must remain inside {repository}: {path}")
    current = repository
    for part in lexical.relative_to(repository).parts:
        current /= part
        if current.is_symlink():
            raise RuntimeError(f"{label} contains a symlink component: {current}")
    resolved = lexical.resolve()
    if repository not in resolved.parents or not resolved.is_dir():
        raise RuntimeError(f"{label} is missing or resolves outside the repository: {path}")
    if resolved not in _VALIDATED_SOURCE_TREES:
        for member in resolved.rglob("*"):
            if member.is_symlink():
                raise RuntimeError(f"{label} contains a symlink: {member}")
            if not member.is_dir() and not member.is_file():
                raise RuntimeError(f"{label} contains a special file: {member}")
        _VALIDATED_SOURCE_TREES.add(resolved)
    return resolved


def validate_project_source_file(path: Path, label: str) -> Path:
    repository = Path(ROOT).resolve()
    lexical = Path(os.path.abspath(path))
    if lexical == repository or repository not in lexical.parents:
        raise RuntimeError(f"{label} must remain inside {repository}: {path}")
    current = repository
    for part in lexical.relative_to(repository).parts:
        current /= part
        if current.is_symlink():
            raise RuntimeError(f"{label} contains a symlink component: {current}")
    resolved = lexical.resolve()
    if repository not in resolved.parents or not resolved.is_file():
        raise RuntimeError(f"{label} is missing or resolves outside the repository: {path}")
    return resolved


def assert_generated_targets_safe(output: Path) -> None:
    allowed = {"README.md", "tasks", "world-image", "lab-agent-image", "dataset"}
    unexpected = sorted(path.name for path in output.iterdir() if path.name not in allowed)
    if unexpected:
        raise RuntimeError(f"generated output root contains unexpected entries: {unexpected}")
    for name in sorted(allowed):
        target = output / name
        if target.is_symlink():
            raise RuntimeError(f"refusing symlinked generated target: {target}")


def remove_generated_directory(path: Path, label: str) -> None:
    """Remove a known generated subtree without leaving a stale empty shell."""
    if path.is_symlink():
        raise RuntimeError(f"refusing to replace symlinked {label}: {path}")
    if path.exists():
        if not path.is_dir():
            raise RuntimeError(f"refusing to replace non-directory {label}: {path}")
        shutil.rmtree(path)


def validate_task_layout(tasks: list[dict]) -> None:
    task_ids: set[str] = set()
    for task in tasks:
        task_id = safe_component(task.get("task_id"), "task_id")
        if task_id in task_ids:
            raise RuntimeError(f"duplicate task_id: {task_id}")
        task_ids.add(task_id)
        phase_names: set[str] = set()
        for phase in (task.get("multi_step") or {}).get("phases") or []:
            name = safe_component(phase.get("name"), f"{task_id} phase name")
            if name in phase_names:
                raise RuntimeError(f"{task_id}: duplicate phase name: {name}")
            phase_names.add(name)
        file_lane = task.get("file_lane") or {}
        validated_deliverables(task)
        skills = file_lane.get("skills")
        if skills is not None:
            if not isinstance(skills, list):
                raise RuntimeError(f"{task_id}: file-lane skills must be a list")
            clean_skills = [
                safe_component(name, f"{task_id} skill name") for name in skills
            ]
            if len(clean_skills) != len(set(clean_skills)):
                raise RuntimeError(f"{task_id}: duplicate file-lane skill name")


def toml_str(s: str) -> str:
    """JSON string escaping is a valid TOML basic string."""
    return json.dumps(str(s))


def load_world(path: str) -> dict:
    with open(path, encoding="utf-8") as f:
        raw = json.load(f)
    return raw.get("world", raw)


def contracts_for_world(world_path: str) -> str:
    """Select the matching frozen contract suite unless explicitly overridden."""
    version = load_world(world_path).get("version") or 0
    suite = "v5" if int(version) >= 21 else ("v4" if int(version) >= 20 else "v3")
    return os.path.join(ROOT, "mcp", suite, "contracts")


def contract_tool_count(contracts: str) -> int:
    total = 0
    for name in sorted(os.listdir(contracts)):
        if not name.endswith(".json"):
            continue
        with open(os.path.join(contracts, name), encoding="utf-8") as f:
            total += sum(
                tool.get("agent_visible") is not False
                for tool in (json.load(f).get("tools") or [])
            )
    return total


# ---------------------------------------------------------------------------
# Per-task files
# ---------------------------------------------------------------------------

def instruction_md(
    task: dict,
    phase: dict | None = None,
    *,
    include_file_deliverables: bool = True,
) -> str:
    prompt = ((phase or {}).get("instruction") or task.get("prompt") or "").strip()
    parts = [prompt]
    # Harbor-native multi-step tasks deliver one instruction per step.  The
    # flattened addendum remains only for single-instruction harnesses.
    session = [] if phase is not None else (task.get("session") or [])
    followups = [s.get("user_text") for s in session if s.get("user_text")]
    if followups:
        parts.append(
            "\n## Follow-up messages\n\n"
            "In the original multi-turn session the user sends these additional "
            "messages, in order, while the work is underway. Treat them as part "
            "of the request:\n\n"
            + "\n".join(f"{i}. {json.dumps(t)}" for i, t in enumerate(followups, 1))
        )
    file_lane = task.get("file_lane") or {}
    if file_lane and not include_file_deliverables:
        parts.append("""
## Input-document lane

The source documents are mounted read-only at `/workspace/documents`. Read the
documents relevant to this step before acting. Final filesystem deliverables
are requested and graded only in the final step; do not anticipate later
instructions.
""")
    elif file_lane:
        deliverables = file_lane.get("deliverables") or []
        if file_lane.get("inputs_only") and not deliverables:
            parts.append("""
## Input-document lane

The source documents are mounted read-only at `/workspace/documents`. Read all
of them before acting. This is an input-only evidence pack: do not create a
separate filesystem deliverable. Complete the requested system-of-record
actions with the product tools; those actions and their evidence trace are
graded deterministically.
""")
        else:
            rendered = "\n".join(f"- `/workspace/output/{name}`" for name in deliverables)
            parts.append(f"""
## File deliverable lane (Harvey LAB parity)

The source documents are mounted read-only at `/workspace/documents`. Write
the required deliverables under `/workspace/output`; expected files:

{rendered or "- Follow the exact output filename stated in the assignment."}

The file lane and the law-firm system-of-record lane are verified separately.
Creating a file does not substitute for filing the work through the product
tools, and filing a record does not substitute for creating the requested
file.
""")
    parts.append(f"""
## Environment

You are an agent operating inside a fully synthetic litigation/corporate
law-firm simulation world ("Eve Litigation" — SIMULATED; no real entities,
clients, or matters). The firm's systems of record are exposed as tools from a
world server on the container network.

Two equivalent ways to use the tools:

1. The `tool` CLI (available on PATH):

   ```bash
   tool list                        # every tool: name + description
   tool schema <name>               # input schema for one tool
   tool call <name> '<json-args>'   # e.g. tool call documents_download '{{"id": 12}}'
   ```

2. MCP (streamable-http, JSON-RPC `tools/list` / `tools/call`) at
   `http://world:8972/mcp` — also declared to MCP-capable harnesses as the
   `lawfirm` server.

Rules:

- Complete the task using the tools. Be precise with record ids and values.
- Read input documents in full before drafting deliverables from them.
- Deliverables must be produced in the systems of record via tools (for
  example `documents_create`); work only described in chat text does not
  count as done.
- Transient tool errors (`rate_limited`, `stale_reference`) are recoverable —
  retry the same call. Ambiguous write acknowledgements ("queued for
  processing") may still have applied: re-read the record instead of writing
  twice.
- When a list/search response exposes a next-page token or URL, continue until
  the relevant result set has been exhausted.

When the work is complete, finish your session; grading is automatic.
""")
    return "\n".join(parts).strip() + "\n"


def task_toml(task: dict, image_tag: str, world_version) -> str:
    tid = task["task_id"]
    prov = task.get("provenance") or {}
    goal = (task.get("goal") or "").strip().replace("\n", " ")
    if len(goal) > 300:
        goal = goal[:297] + "..."
    keywords = [k for k in [
        task.get("difficulty_tier"), task.get("complexity"),
        (prov.get("source_workflow") or "").split("/")[-1] or None,
    ] if k]
    multi_step = task.get("multi_step") or {}
    lines = [
        'schema_version = "1.4"',
    ]
    if multi_step:
        lines.append(f'multi_step_reward_strategy = {toml_str(multi_step.get("reward_strategy") or "mean")}')
    lines += ["",
        "[task]",
        f'name = "legal-agent-simulation/{tid.replace("_", "-")}"',
        f'version = "{world_version}.0.0"',
        f"description = {toml_str(goal or tid)}",
        "authors = []",
        f"keywords = {json.dumps(keywords)}",
        "",
        "[metadata]",
        f"task_id = {toml_str(tid)}",
        f"difficulty = {toml_str(task.get('difficulty_tier') or '')}",
        f"complexity = {toml_str(task.get('complexity') or '')}",
        f"acceptance_label = {toml_str(task.get('acceptance_label') or '')}",
        f"source_workflow = {toml_str(prov.get('source_workflow') or '')}",
        f"method = {toml_str(task.get('method') or '')}",
        f"world_image = {toml_str(image_tag)}",
        f"disclaimer = {toml_str(DISCLAIMER)}",
        f"file_lane = {str(bool(task.get('file_lane'))).lower()}",
        f"lab_source_task = {toml_str((task.get('file_lane') or {}).get('source_task') or '')}",
        f"lab_source_commit = {toml_str((task.get('file_lane') or {}).get('source_commit') or '')}",
        "",
        "[verifier]",
        "timeout_sec = 180.0",
        "",
        "[agent]",
        "timeout_sec = 1800.0",
        "",
        "[environment]",
        "build_timeout_sec = 900.0",
        "cpus = 1",
        "memory_mb = 2048",
        "storage_mb = 10240",
        "gpus = 0",
        "",
        "[[environment.mcp_servers]]",
        'name = "lawfirm"',
        'transport = "streamable-http"',
        'url = "http://world:8972/mcp"',
        "",
    ]
    for phase in multi_step.get("phases") or []:
        lines.extend([
            "[[steps]]",
            f'name = {toml_str(phase["name"])}',
            f'min_reward = {float(phase.get("min_reward", 1.0))}',
            "",
        ])
    return "\n".join(lines)


AGENT_DOCKERFILE = """\
# Agent container. The world (tools, verifiers, state) lives in the separate
# `world` compose service — see docker-compose.yaml.
FROM python:3.12-slim@sha256:229a2c5bfa27522db7815ea81f9bed70af17ccb9de9fc7ad142b1877b5830d36
COPY tool /usr/local/bin/tool
RUN chmod +x /usr/local/bin/tool
ENV LAWFIRM_MCP=http://world:8972/mcp
WORKDIR /app
"""


def lab_agent_dockerfile(base_image: str) -> str:
    return f"""\
# File-lane agent container. The base is the dependency-locked production
# derivative of Harvey LAB and contains LibreOffice, pandoc, and format parsers.
FROM {base_image}
COPY tool /usr/local/bin/tool
RUN chmod +x /usr/local/bin/tool
COPY skills /workspace/skills
ENV LAWFIRM_MCP=http://world:8972/mcp \\
    WORKSPACE_DIR=/workspace \\
    DOCUMENTS_DIR=/workspace/documents \\
    OUTPUT_DIR=/workspace/output
RUN mkdir -p /workspace/documents /workspace/output
WORKDIR /workspace
"""


def compose_yaml(task_id: str, image_tag: str, file_lane: bool = False) -> str:
    documents_mount = """
    volumes:
      - type: bind
        source: ./documents
        target: /workspace/documents
        read_only: true
""" if file_lane else ""
    return f"""\
# Merged on top of Harbor's base compose config; `main` (the agent container)
# is configured by Harbor automatically. The shared world image is built once:
#   python3 harbor/generate.py --build-image
services:
  main:
    depends_on:
      world:
        condition: service_healthy
{documents_mount.rstrip()}

  world:
    image: {image_tag}
    environment:
      TASK_ID: {task_id}
    expose:
      - "8972"
    healthcheck:
      test: ["CMD", "python3", "-c", "import urllib.request; urllib.request.urlopen('http://127.0.0.1:8972/health', timeout=2)"]
      interval: 2s
      timeout: 5s
      retries: 60
      start_period: 5s
"""


def validated_deliverables(task: dict) -> list[str]:
    values = (task.get("file_lane") or {}).get("deliverables") or []
    clean = []
    for value in values:
        if not isinstance(value, str) or "\\" in value:
            raise RuntimeError(f"{task['task_id']}: unsafe deliverable path {value!r}")
        path = PurePosixPath(value)
        if (
            path.is_absolute()
            or ".." in path.parts
            or not path.name
            or path.as_posix() != value
        ):
            raise RuntimeError(f"{task['task_id']}: unsafe deliverable path {value!r}")
        clean.append(path.as_posix())
    return clean


def test_sh(
    task: dict,
    phase: str | None = None,
    *,
    include_file_deliverables: bool = True,
) -> str:
    file_lane = (task.get("file_lane") or {}) if include_file_deliverables else {}
    expected = json.dumps(validated_deliverables(task) if file_lane else [])
    file_assertions = json.dumps(file_lane.get("assertions") or {})
    verify_body = json.dumps({"phase": phase} if phase else {})
    return f"""\
#!/bin/bash
# Verifier: ask the world container for the trial verdict (shipped VCode,
# executed against the session's final state + the recorded tool trace),
# then emit the Harbor reward file.
python3 - <<'PYEOF'
import hashlib, json, os, re, shutil, unicodedata, urllib.request, zipfile
from xml.etree import ElementTree as ET

expected = {expected}
file_assertions = {file_assertions}
logs_root = os.environ.get("HARBOR_LOGS", "/logs")
artifact_root = os.path.join(logs_root, "artifacts")
verifier_root = os.path.join(logs_root, "verifier")
output_root = os.environ.get("WORKSPACE_OUTPUT", "/workspace/output")
verify_url = os.environ.get("WORLD_VERIFY_URL", "http://world:8972/verify")
os.makedirs(artifact_root, exist_ok=True)
artifacts = []
rejected_artifacts = []
if os.path.isdir(output_root):
    for root, _, files in os.walk(output_root):
        for name in sorted(files):
            source = os.path.join(root, name)
            relative = os.path.relpath(source, output_root)
            if os.path.islink(source):
                rejected_artifacts.append({{"path": relative, "reason": "symlink"}})
                continue
            target = os.path.join(artifact_root, relative)
            os.makedirs(os.path.dirname(target), exist_ok=True)
            shutil.copy2(source, target)
            artifacts.append({{"path": relative, "bytes": os.path.getsize(source)}})
output_real = os.path.realpath(output_root)
def valid_expected_file(name):
    candidate = os.path.join(output_root, name)
    real = os.path.realpath(candidate)
    return (
        os.path.commonpath([output_real, real]) == output_real and
        not os.path.islink(candidate) and os.path.isfile(candidate) and
        os.path.getsize(candidate) > 0
    )
file_contract_passed = None if not expected else all(valid_expected_file(name) for name in expected)

def normalize(value):
    value = unicodedata.normalize("NFKC", str(value or "")).casefold()
    value = value.replace("\\u00a0", " ").replace("–", "-").replace("—", "-")
    return re.sub(r"\\s+", " ", value).strip()

def has_anchor(text, anchor):
    needle = normalize(anchor)
    if not needle:
        return False
    left = r"(?<![\\w.,])" if needle[0].isdigit() else r"(?<![\\w])"
    right = r"(?![\\w]|\\.\\d)" if needle[-1].isdigit() else r"(?![\\w])"
    return bool(re.search(left + re.escape(needle) + right, text))

def output_text(path):
    suffix = os.path.splitext(path)[1].lower()
    if suffix in (".txt", ".md", ".json"):
        return open(path, encoding="utf-8", errors="replace").read()
    if suffix == ".xlsx":
        import openpyxl
        workbook = openpyxl.load_workbook(path, read_only=True, data_only=False)
        values = []
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                values.extend(str(value) for value in row if value is not None)
        workbook.close()
        return "\\n".join(values)
    if suffix in (".docx", ".pptx"):
        values = []
        with zipfile.ZipFile(path) as archive:
            prefix = "word/" if suffix == ".docx" else "ppt/"
            for name in sorted(archive.namelist()):
                if name.startswith(prefix) and name.endswith(".xml"):
                    root = ET.fromstring(archive.read(name))
                    values.extend(node.text or "" for node in root.iter()
                                  if node.tag.rsplit("}}", 1)[-1] == "t")
        return "\\n".join(values)
    if suffix == ".pdf":
        from pypdf import PdfReader
        return "\\n".join(page.extract_text() or "" for page in PdfReader(path).pages)
    return open(path, "rb").read().decode("utf-8", errors="replace")

file_criterion_results = []
file_parse_errors = []
text_by_file = {{}}
if file_contract_passed:
    criteria_to_check = []
    if isinstance(file_assertions, dict):
        for name, criteria in file_assertions.items():
            criteria_to_check.extend({{**criterion, "deliverables": [name]}} for criterion in criteria)
    else:
        criteria_to_check = file_assertions
    names_to_parse = set(expected) | {{name for criterion in criteria_to_check
                                      for name in criterion.get("deliverables") or []}}
    for name in sorted(names_to_parse):
        try:
            text_by_file[name] = normalize(output_text(os.path.join(output_root, name)))
        except Exception as error:
            file_parse_errors.append({{"path": name, "error": repr(error)}})
            text_by_file[name] = ""
    for criterion in criteria_to_check:
        targets = criterion.get("deliverables") or expected
        text = "\\n".join(text_by_file.get(name, "") for name in targets)
        missing = []
        for group in criterion.get("anchor_groups") or []:
            if not any(has_anchor(text, anchor) for anchor in group):
                missing.append(group)
        file_criterion_results.append({{
            "criterion_id": criterion.get("criterion_id"),
            "deliverables": targets,
            "passed": not missing,
            "missing_anchor_groups": missing,
        }})
file_content_passed = (None if not file_assertions else
                       bool(file_criterion_results) and
                       not file_parse_errors and
                       all(row["passed"] for row in file_criterion_results))
file_passed = (file_contract_passed if file_content_passed is None else
               bool(file_contract_passed and file_content_passed))

verdict, out = None, {{"reward": 0.0, "passed": 0.0}}
try:
    req = urllib.request.Request(verify_url, method="POST")
    req.add_header("Content-Type", "application/json")
    with urllib.request.urlopen(req, data={verify_body!r}.encode(), timeout=150) as res:
        verdict = json.loads(res.read().decode() or "{{}}")
    out["reward"] = round(float(verdict.get("reward") or 0.0), 4)
    out["passed"] = 1.0 if verdict.get("passed") else 0.0
except Exception as e:  # world unreachable / verifier crash -> reward 0
    verdict = {{"error": repr(e)}}

state_passed = bool(out["passed"])
state_digests = (verdict or {{}}).get("filed_text_sha256") or {{}}
file_digests = {{name: hashlib.sha256(normalize(text_by_file.get(name, "")).encode()).hexdigest()
                for name in expected if name in text_by_file}}
cross_lane_match = (None if not expected or not state_digests else
                    all(file_digests.get(name) == state_digests.get(name) for name in expected))
lane = {{
    "enabled": bool(expected),
    "grade_kind": "determinate" if file_assertions else "output_contract_only",
    "expected": expected,
    "artifacts": artifacts,
    "rejected_artifacts": rejected_artifacts,
    "file_contract_passed": file_contract_passed,
    "file_content_passed": file_content_passed,
    "file_criterion_results": file_criterion_results,
    "file_parse_errors": file_parse_errors,
    "file_passed": file_passed,
    "state_passed": state_passed,
    "cross_lane_match": cross_lane_match,
    "file_text_sha256": file_digests,
    "state_text_sha256": state_digests,
    "lane_split": None if file_passed is None else file_passed != state_passed,
}}
# Harbor's VerifierResult contract accepts numeric reward channels only.  Keep
# the nullable/boolean diagnostic values in file-lane.json, and expose each
# applicable channel here as 0.0/1.0 so a real Harbor run can ingest it.
out["state_passed"] = 1.0 if state_passed else 0.0
if file_passed is not None:
    out["file_passed"] = 1.0 if file_passed else 0.0
    out["lane_split"] = 1.0 if lane["lane_split"] else 0.0
if cross_lane_match is not None:
    out["cross_lane_match"] = 1.0 if cross_lane_match else 0.0

os.makedirs(verifier_root, exist_ok=True)
with open(os.path.join(verifier_root, "verdict.json"), "w") as f:
    json.dump(verdict, f, indent=1)
with open(os.path.join(verifier_root, "reward.json"), "w") as f:
    json.dump(out, f)
with open(os.path.join(verifier_root, "file-lane.json"), "w") as f:
    json.dump(lane, f, indent=1)
print(json.dumps({{"passed": out["passed"], "reward": out["reward"],
                  "file_passed": file_passed, "lane_split": lane["lane_split"],
                  "failed_conditions": (verdict or {{}}).get("failed_conditions")}}))
PYEOF
"""


def oracle_file_outputs(task: dict) -> dict[str, str]:
    """Build deterministic file-lane oracle text from compiled anchor groups."""
    outputs: dict[str, list[str]] = {name: [] for name in validated_deliverables(task)}
    for criterion in (task.get("file_lane") or {}).get("assertions") or []:
        targets = criterion.get("deliverables") or list(outputs)
        fragment = " | ".join(str(group[0]) for group in criterion.get("anchor_groups") or [] if group)
        if not fragment:
            continue
        for name in targets:
            if name in outputs:
                outputs[name].append(fragment)
    return {name: "\n".join(lines) or "Completed deliverable."
            for name, lines in outputs.items()}


def solve_sh(
    token: str,
    task: dict | None = None,
    phase: str | None = None,
    *,
    include_file_deliverables: bool = True,
) -> str:
    import base64
    file_task = (task or {}) if include_file_deliverables else {}
    payload = base64.b64encode(json.dumps(oracle_file_outputs(file_task)).encode()).decode()
    file_writer = f"""\
python3 - <<'PYEOF'
import base64, json
from pathlib import Path

outputs = json.loads(base64.b64decode({payload!r}).decode())
root = Path('/workspace/output')
root.mkdir(parents=True, exist_ok=True)
for name, body in outputs.items():
    target = root / name
    target.parent.mkdir(parents=True, exist_ok=True)
    suffix = target.suffix.casefold()
    if suffix == '.docx':
        from docx import Document
        document = Document()
        for line in body.splitlines() or ['Completed deliverable.']:
            document.add_paragraph(line)
        document.save(target)
    elif suffix == '.xlsx':
        from openpyxl import Workbook
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = 'Determinate Findings'
        for row, line in enumerate(body.splitlines() or ['Completed deliverable.'], 1):
            sheet.cell(row=row, column=1, value=line)
        workbook.save(target)
    elif suffix == '.pptx':
        from pptx import Presentation
        from pptx.util import Inches
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(0.75), Inches(0.75), Inches(8.5), Inches(5.5))
        box.text_frame.text = body or 'Completed deliverable.'
        presentation.save(target)
    else:
        target.write_text(body + '\\n', encoding='utf-8')
PYEOF
""" if payload != "e30=" else ""
    solve_body = json.dumps({"phase": phase} if phase else {})
    return f"""\
#!/bin/bash
# Oracle solution: the world container replays this task's reference walk
# through the live session. The raw proof token exists only in this solution
# file; the world image contains its SHA-256 digest and this file is never
# present during ordinary agent runs.
set -e
{file_writer}python3 - <<'PYEOF'
import urllib.request

request = urllib.request.Request('http://world:8972/solve', method='POST')
request.add_header('X-Solve-Token', {token!r})
request.add_header('Content-Type', 'application/json')
with urllib.request.urlopen(request, data={solve_body!r}.encode(), timeout=150) as response:
    print(response.read().decode())
PYEOF
"""


def root_readme(
    task_count: int,
    world_path: str | os.PathLike[str],
    contracts_dir: str | os.PathLike[str],
    image_tag: str,
    world: dict,
    runtime_tool_count: int,
) -> str:
    """Render the byte-exact top-level export documentation."""
    return f"""\
# legal-agent-simulation — Harbor tasks

{task_count} Harbor-format tasks generated from `{os.path.relpath(world_path, ROOT)}`
(one per world task; regenerate with `python3 harbor/generate.py`).

Contract suite: `{os.path.relpath(contracts_dir, ROOT)}`.

{DISCLAIMER}

## One-time setup

Build the shared world image (world runtime + world doc + product contracts):

```bash
python3 harbor/generate.py --build-image        # tags {image_tag}
```

## Run

```bash
harbor run -p "dist/harbor/tasks/task_003" -a claude-code -m anthropic/claude-sonnet-5
harbor run -p "dist/harbor/tasks/task_003" -a oracle       # reference-walk sanity check
```

Multi-container tasks require Harbor's **docker** environment provider
(compose networking); cloud providers are not supported for these tasks.

## Architecture

- `world` service (shared image `{image_tag}`): the executable law-firm
  world — {len(world["tables"])} product-system tables hydrated from the world
  doc, {runtime_tool_count} contract-defined tools and zero synthesized
  name-family tools,
  deterministic seeded friction, and shipped VCode verifiers.
  A per-trial shim creates the task's session, records the tool trace, and
  exposes `POST /mcp` (JSON-RPC), `POST /verify`, and token-gated `POST /solve`.
- `main` (agent container): pinned python + the `tool` CLI. It contains no
  world document, verifier code, or reference walks.
- `tests/test.sh` fetches the VCode verdict and writes `reward.json`
  (`reward` = graded fraction with anti-hack vetoes, `passed` = strict bool).
- `solution/solve.sh` triggers the oracle reference walk server-side — the
  same walk `world/local/oracle.py` proves against all {len(world['tasks'])}
  tasks in this generated world.
"""


# ---------------------------------------------------------------------------
# Assembly
# ---------------------------------------------------------------------------

def write(path: str, content: str, executable: bool = False) -> None:
    destination = Path(path)
    if destination.is_symlink():
        raise RuntimeError(f"refusing to overwrite symlink: {destination}")
    destination.parent.mkdir(parents=True, exist_ok=True)
    try:
        with destination.open("x", encoding="utf-8", newline="\n") as handle:
            handle.write(content)
    except FileExistsError:
        pass
    else:
        destination.chmod(0o755 if executable else 0o644)
        return
    temporary_path: Path | None = None
    try:
        with tempfile.NamedTemporaryFile(
            mode="w",
            encoding="utf-8",
            newline="\n",
            prefix=f".{destination.name}.",
            dir=destination.parent,
            delete=False,
        ) as handle:
            temporary_path = Path(handle.name)
            handle.write(content)
        temporary_path.chmod(0o755 if executable else 0o644)
        os.replace(temporary_path, destination)
        temporary_path = None
    finally:
        if temporary_path is not None:
            temporary_path.unlink(missing_ok=True)


def link_or_copy(source: str, destination: str) -> str:
    """Copy-tree callback that avoids duplicating the multi-GB local corpus."""
    try:
        os.link(source, destination)
    except OSError:
        shutil.copy2(source, destination)
    return destination


def stage_file_lane(task: dict, task_dir: str) -> None:
    config = task.get("file_lane") or {}
    source = config.get("documents_source")
    if not source:
        raise RuntimeError(f"{task['task_id']}: file_lane.documents_source is required")
    source_path = Path(source)
    if not source_path.is_absolute():
        source_path = Path(ROOT) / source_path
    source_path = validate_project_source_tree(
        source_path, f"{task['task_id']} file-lane documents"
    )
    destination = Path(task_dir) / "environment" / "documents"
    shutil.copytree(source_path, destination, copy_function=link_or_copy)

    configured_skills = config.get("skills_source")
    if configured_skills:
        skills_source = Path(configured_skills)
    else:
        skills_source = (Path(ROOT) / "research" / "repos" / "harveyai@harvey-labs" /
                         "harness" / "skills")
        if not skills_source.is_dir():
            skills_source = Path(ROOT) / "research" / "harvey-recovery" / "skills"
    if not skills_source.is_absolute():
        skills_source = Path(ROOT) / skills_source
    skills_source = validate_project_source_tree(
        skills_source, f"{task['task_id']} skills source"
    )
    skills_destination = Path(task_dir) / "environment" / "skills"
    requested = config["skills"] if "skills" in config else ["docx", "xlsx", "pptx"]
    requested = [safe_component(name, f"{task['task_id']} skill name") for name in requested]
    skills_destination.mkdir(parents=True, exist_ok=True)
    for name in requested:
        source_skill = skills_source / name
        if not (source_skill / "SKILL.md").is_file():
            raise RuntimeError(f"{task['task_id']}: Harvey LAB skill missing: {source_skill}")
        shutil.copytree(source_skill, skills_destination / name, copy_function=link_or_copy)


def assemble_world_image(out: str, world_path: str, contracts_dir: str | None = None) -> str:
    """Copy runtime + world doc + shim into the shared image build context."""
    img_path = Path(out) / "world-image"
    reset_generated_directory(img_path, "world image context")
    img = str(img_path)
    local = os.path.join(ROOT, "world", "local")
    for name in WORLD_RUNTIME_FILES:
        shutil.copy2(os.path.join(local, name), os.path.join(img, name))
    for name in WORLD_IMAGE_TEMPLATE_FILES:
        shutil.copy2(os.path.join(HERE, "world-image", name),
                     os.path.join(img, name))
    shutil.copy2(world_path, os.path.join(img, "world.json"))
    contracts = os.path.abspath(contracts_dir or contracts_for_world(world_path))
    dst = os.path.join(img, "contracts")
    if os.path.isdir(dst):
        shutil.rmtree(dst)
    shutil.copytree(contracts, dst)
    evidence_destination = Path(img) / "corpus"
    if evidence_destination.exists():
        shutil.rmtree(evidence_destination)
    evidence_destination.mkdir()
    world = load_world(world_path)
    evidence_kinds = sorted({
        str((task.get("evidence_store") or {}).get("kind"))
        for task in world.get("tasks") or [] if task.get("evidence_store")
    })
    for kind in evidence_kinds:
        if kind not in {"lab", "ch"}:
            raise RuntimeError(f"unsupported packaged evidence kind: {kind}")
        source_index = Path(ROOT) / "world" / "corpus" / kind / "index.sqlite"
        if source_index.is_symlink() or not source_index.is_file():
            raise RuntimeError(f"{kind} evidence index missing: {source_index}")
        target = evidence_destination / kind
        target.mkdir()
        link_or_copy(str(source_index), str(target / "index.sqlite"))
    return img


def assemble_lab_agent_image(out: str) -> str:
    """Assemble the locked production derivative of the exact Harvey sandbox."""
    upstream = Path(ROOT) / "research" / "harvey-recovery" / "sandbox"
    upstream = validate_project_source_tree(upstream, "Harvey LAB sandbox source")
    if not (upstream / "Dockerfile").is_file():
        raise RuntimeError(f"tracked Harvey LAB sandbox source missing: {upstream}")
    template = Path(HERE) / "lab-agent-image"
    template = validate_project_source_tree(template, "LAB image template")
    required_template_files = (
        "Dockerfile",
        "debian.sources",
        "requirements.in",
        "requirements.lock",
        "package.json",
        "package-lock.json",
    )
    for name in required_template_files:
        if not (template / name).is_file():
            raise RuntimeError(f"locked LAB image input missing: {template / name}")

    destination = Path(out) / "lab-agent-image"
    if destination.is_symlink():
        raise RuntimeError(f"refusing symlinked LAB image context: {destination}")
    if destination.exists():
        if not destination.is_dir():
            raise RuntimeError(f"refusing non-directory LAB image context: {destination}")
        shutil.rmtree(destination)
    # This context is tiny. Use independent copies so replacing the staged
    # Dockerfile cannot mutate the exact recovery source through a hard link.
    shutil.copytree(upstream, destination)
    for name in required_template_files:
        shutil.copy2(template / name, destination / name)
    return str(destination)


def resolve_solve_token(token_path: str) -> str:
    """Use the release secret, an existing export token, or a new local token."""
    configured = os.environ.get("HARBOR_SOLVE_TOKEN", "").strip()
    if configured:
        token = configured
    elif Path(token_path).is_symlink():
        raise RuntimeError(f"refusing symlinked solve token: {token_path}")
    elif os.path.isfile(token_path):
        token = Path(token_path).read_text("utf-8").strip()
    else:
        token = secrets.token_hex(16)
    if not (32 <= len(token) <= 128 and all(character in "0123456789abcdef" for character in token)):
        raise RuntimeError(
            "HARBOR_SOLVE_TOKEN/export solve token must be 32-128 lowercase hex characters"
        )
    return token


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--world", default=os.path.join(ROOT, "world", "blobfish",
                                                    "world-v21.json"))
    ap.add_argument("--contracts", default="",
                    help="product-contract directory; defaults to mcp/v5 for world v21+, "
                         "mcp/v4 for v20, and mcp/v3 for historical worlds")
    ap.add_argument("--out", default=os.path.join(ROOT, "dist", "harbor"))
    ap.add_argument("--tasks", default="", help="comma-separated task_id filter")
    ap.add_argument("--image-tag",
                    default="ghcr.io/blobfishai/legal-agent-sim-world:v21",
                    help="world image reference baked into every task's compose "
                         "file; --build-image tags the local build with it")
    ap.add_argument("--build-image", action="store_true")
    ap.add_argument("--lab-agent-image",
                    default="ghcr.io/blobfishai/legal-agent-sim-agent-lab:v21")
    ap.add_argument("--build-lab-agent-image", action="store_true",
                    help="build the heavy file-lane base from the pinned LAB sandbox")
    args = ap.parse_args()

    world_path = validate_project_source_file(Path(args.world), "world source")
    world = load_world(str(world_path))
    validate_task_layout(world["tasks"])
    contracts_dir = validate_project_source_tree(
        Path(args.contracts or contracts_for_world(str(world_path))),
        "contract source",
    )
    runtime_tool_count = contract_tool_count(str(contracts_dir))
    wanted = {t for t in args.tasks.split(",") if t}
    tasks = [t for t in world["tasks"] if not wanted or t["task_id"] in wanted]
    if wanted and len(tasks) != len(wanted):
        sys.exit(f"unknown task ids: {sorted(wanted - {t['task_id'] for t in tasks})}")

    out_path = resolve_output_root(args.out)
    out_path.mkdir(parents=True, exist_ok=True)
    assert_generated_targets_safe(out_path)
    # A dataset is derived from the freshly generated task bytes by the locked
    # production wrapper. Never leave a manifest from the previous generation
    # beside new task packages.
    remove_generated_directory(out_path / "dataset", "Harbor dataset")
    out = str(out_path)

    # The encrypted release secret lets independently generated full-task and
    # production-image exports share one hidden oracle credential. Without it,
    # local regeneration remains stable by retaining the ignored token file.
    token_path = os.path.join(out, "world-image", "solve-token.txt")
    token = resolve_solve_token(token_path)
    img_dir = assemble_world_image(out, str(world_path), str(contracts_dir))
    lab_img_dir = assemble_lab_agent_image(out)
    write(token_path, token + "\n")

    if args.build_lab_agent_image:
        command = ["docker", "build", "-t", args.lab_agent_image, lab_img_dir]
        print("+", " ".join(command))
        subprocess.run(command, check=True)

    tool_src = Path(HERE, "agent-image", "tool").read_text("utf-8")
    tasks_root = os.path.join(out, "tasks")
    reset_generated_directory(Path(tasks_root), "Harbor tasks tree")
    for task in tasks:
        tid = task["task_id"]
        d = os.path.join(tasks_root, tid)
        multi_step = task.get("multi_step") or {}
        phases = multi_step.get("phases") or []
        if not phases:
            write(os.path.join(d, "instruction.md"), instruction_md(task))
        write(os.path.join(d, "task.toml"),
              task_toml(task, args.image_tag, world.get("version")))
        if task.get("file_lane"):
            stage_file_lane(task, d)
            dockerfile = lab_agent_dockerfile(args.lab_agent_image)
        else:
            dockerfile = AGENT_DOCKERFILE
        write(os.path.join(d, "environment", "Dockerfile"), dockerfile)
        write(os.path.join(d, "environment", "tool"), tool_src, executable=True)
        write(os.path.join(d, "environment", "docker-compose.yaml"),
              compose_yaml(tid, args.image_tag, bool(task.get("file_lane"))))
        if phases:
            for index, phase in enumerate(phases):
                include_file_deliverables = index == len(phases) - 1
                step_dir = os.path.join(d, "steps", phase["name"])
                write(
                    os.path.join(step_dir, "instruction.md"),
                    instruction_md(
                        task,
                        phase,
                        include_file_deliverables=include_file_deliverables,
                    ),
                )
                write(os.path.join(step_dir, "tests", "test.sh"),
                      test_sh(
                          task,
                          phase["name"],
                          include_file_deliverables=include_file_deliverables,
                      ), executable=True)
                write(os.path.join(step_dir, "solution", "solve.sh"),
                      solve_sh(
                          token,
                          task,
                          phase["name"],
                          include_file_deliverables=include_file_deliverables,
                      ), executable=True)
        else:
            write(os.path.join(d, "tests", "test.sh"), test_sh(task), executable=True)
            write(os.path.join(d, "solution", "solve.sh"), solve_sh(token, task),
                  executable=True)

    write(
        os.path.join(out, "README.md"),
        root_readme(
            len(tasks),
            world_path,
            contracts_dir,
            args.image_tag,
            world,
            runtime_tool_count,
        ),
    )

    print(f"generated {len(tasks)} Harbor tasks -> {tasks_root}")
    print(f"world image context -> {img_dir}")

    if args.build_image:
        token_hash = hashlib.sha256(token.encode()).hexdigest()
        cmd = ["docker", "build", "-t", args.image_tag,
               "--build-arg", f"ORACLE_PROOF_SHA256={token_hash}", img_dir]
        print("+", " ".join(cmd))
        subprocess.run(cmd, check=True)


if __name__ == "__main__":
    main()
